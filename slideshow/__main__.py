from tqdm import tqdm
import pptx
from io import BytesIO
import whatimage
import pyheif
from PIL import Image
import os
import cv2
import sys
from pkg_resources import resource_filename


def detectFaces(gray,loadedCascade):
    for orientation in range(4): # 4 shots
        # Detect faces in the image
        faces = loadedCascade.detectMultiScale(
            gray,
            scaleFactor=1.1,
            minNeighbors=5,
            minSize=(30, 30),
            flags = cv2.CASCADE_SCALE_IMAGE
        )
        if len(faces) > 0:
            return orientation, True
        # Rotate it for the next time
        gray = cv2.rotate(gray, cv2.ROTATE_90_CLOCKWISE)
    return orientation, False

# Detect which side of the photo is brightest. Hopefully it will be the sky.
def detectBrightest(image):
    image_scale = 4 # This scale factor doesn't matter much. It just gives us less pixels to iterate over later
    newsize = (round(image.shape[1] / image_scale), round(image.shape[0] / image_scale)) # find new size
    small_img = cv2.resize(image, newsize, cv2.INTER_AREA)

    # Take the top 1/3, right 1/3, etc. to compare for brightness
    width = int(newsize[0])
    height = int(newsize[1])
    top = small_img[0: int(height / 3), 0:width]
    right = small_img[0: height, int(width / 3 * 2): width]
    left = small_img[0: height, 0: int(width / 3)]
    bottom = small_img[int(height / 3 * 2): height, 0: height]
    sides = {'top':top,'left':left,'bottom':bottom,'right':right}
    # Find the brightest side
    greatest = 0
    winning = 'top'
    for name in sides:
        sidelum = 0
        side = sides[name]
        for x in range(side.shape[0] - 1):
            for y in range(side.shape[1] - 1):
                sidelum = sidelum + side[x,y]
        sidelum = sidelum / (side.shape[0] * side.shape[1])
        if sidelum > greatest:
            winning = name
        # Output based on # of 90 degree clockwise rotations
        returns = {'top':0,
                   'left':90,
                   'bottom':180,
                   'right':270}
        # return the winner
        return (360 - returns[winning]) / 90

def get_haar_cascade(name):
    if os.path.exists(os.path.join("haarcascades", name)):
        path = os.path.join("haarcascades", name)
    else:
        path = resource_filename("slideshow", "haarcascades/" + name)
    return cv2.CascadeClassifier(path)

# Get more at: https://code.ros.org/svn/opencv/tags/latest_tested_snapshot/opencv/data/haarcascades/
# Listed in order most likely to appear in a photo
cascades = list(map(get_haar_cascade, ["haarcascade_frontalface_alt.xml", "haarcascade_profileface.xml", "haarcascade_fullbody.xml"]))

# Try a couple different detection methods
def trydetect(path):
    CV_LOAD_IMAGE_GRAYSCALE = 0
    # Load some things that we'll use during each loop so we don't keep re-creating them
    image = cv2.imread(path)
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    for cascade in cascades:
        orientation, found_faces = detectFaces(gray, cascade)
        if found_faces:
            return orientation
    return detectBrightest(gray) # no faces found, use the brightest side for orientation instead

def decodeHEICImage(bytesIo):
    fmt = whatimage.identify_image(bytesIo.read())
    if fmt in ['heic', 'avif']:
        i = pyheif.read_heif(bytesIo.read())
        # Extract metadata etc

        #for metadata in i.metadata or []:
        #    if metadata['type']=='Exif':
        #        # do whatever

        # Convert to other file format like jpeg
        s = BytesIO()
        pi = Image.frombytes(mode=i.mode, size=i.size, data=i.data)
        pi.save(s, format="jpeg")
        return s

def is_photo(filename: str):
    #photo_types = ["jpg", "png", "heic"]
    photo_types = ["jpg", "png"]
    for typ in photo_types:
        if filename.startswith("."):
            return False
        elif filename.lower().endswith("." + typ):
            return True
    return False

def get_photos_list():
    print("Please input the location of your photos you would like to use")
    photos_dir = input(">> ")
    while not os.path.exists(photos_dir):
        print("Directory not found.  Please try again.")
        photos_dir = input(">> ")
    print("Directory Found.")
    onlyfiles = [os.path.join(photos_dir, f) for f in os.listdir(photos_dir) if is_photo(f)]
    print("Found {} files.".format(len(onlyfiles)))
    return onlyfiles

def get_output_filename():
    print("Please output the location and filename for the final slideshow")
    slideshow_dir = input("(Leave blank for slideshow.pptx in the current directory) >> ")
    return "slideshow.pptx" if slideshow_dir == "" or slideshow_dir is None else slideshow_dir

def get_slide_layout(prs):
    SLD_LAYOUT_TITLE = 0
    SLD_LAYOUT_TITLE_AND_CONTENT = 1
    SLD_LAYOUT_SECTION_HEADER = 2
    SLD_TWO_CONTENT = 3
    SLD_COMPARISON = 4
    SLD_TITLE_ONLY = 5
    SLD_LAYOUT_BLANK = 6
    SLD_CONTENT_WITH_CAPTION = 7
    SLD_PICTURE_WITH_CAPTION = 8
    slide_layout = prs.slide_layouts[SLD_LAYOUT_BLANK]
    return slide_layout

def create_slideshow(photos_list, output_filename):
    prs = pptx.Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    inch_to_emu = 914400

    for photo in tqdm(photos_list):
        slide = prs.slides.add_slide(get_slide_layout(prs))
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = pptx.dml.color.RGBColor(0, 0, 0)

        if photo.endswith("heic"):
            with open(photo, "rb") as raw_photo:
                b = decodeHEICImage(raw_photo)

                width, height = pht.size
                is_wider_than_tall = max(width, height) == width
                is_square = width == height
                aspect_ratio = float(width) / float(height)
                if is_wider_than_tall:
                    new_width = slide_width
                    new_height = float(new_width) / aspect_ratio
                    left = pptx.util.Inches(0)

                else:
                    new_height = slide_height
                    new_width = float(new_height) * aspect_ratio
                    top = pptx.util.Inches(0)

                pic = slide.shapes.add_picture(b, ((slide_width - new_width) / 2), ((slide_height - new_height) / 2), new_width, new_height)

        else:
            orientation = trydetect(photo)
            with Image.open(photo) as pht:
                pht = pht.rotate(360 - (90 * ((orientation) % 3)))
                width, height = pht.size
                is_wider_than_tall = max(width, height) == width
                is_square = width == height
                aspect_ratio = float(width) / float(height)
                if is_wider_than_tall:
                    new_width = slide_width
                    new_height = float(new_width) / aspect_ratio
                    left = pptx.util.Inches(0)
                else:
                    new_height = slide_height
                    new_width = aspect_ratio / float(new_height)
                    top = pptx.util.Inches(0)

                #calculate max width/height for target size
                b = BytesIO()
                if photo.endswith("jpg") or photo.endswith("jpeg"):
                    pht.save(b, 'jpeg')
                elif photo.endswith("png"):
                    pht.save(b, 'png')
                pic = slide.shapes.add_picture(b, ((slide_width - new_width) / 2), ((slide_height - new_height) / 2), new_width, new_height)

    prs.save(output_filename)

if __name__ == "__main__":
    photos_list = get_photos_list()
    output_filename = get_output_filename()
    create_slideshow(photos_list, output_filename)
